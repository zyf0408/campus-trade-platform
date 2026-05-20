from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

COLOR_BG = RGBColor(0xF5, 0xF7, 0xFA)
COLOR_PRIMARY = RGBColor(0x1A, 0x73, 0xE8)
COLOR_PRIMARY_DARK = RGBColor(0x0D, 0x47, 0xA1)
COLOR_ACCENT = RGBColor(0x00, 0xAC, 0xC1)
COLOR_GREEN = RGBColor(0x34, 0xA8, 0x53)
COLOR_ORANGE = RGBColor(0xF9, 0xAB, 0x00)
COLOR_RED = RGBColor(0xEA, 0x43, 0x35)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_DARK = RGBColor(0x20, 0x2A, 0x44)
COLOR_GRAY = RGBColor(0x5F, 0x6B, 0x7A)
COLOR_LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
COLOR_LIGHT_GREEN = RGBColor(0xE6, 0xF4, 0xEA)
COLOR_LIGHT_ORANGE = RGBColor(0xFE, 0xF7, 0xE0)
COLOR_LIGHT_RED = RGBColor(0xFC, 0xE8, 0xE6)
COLOR_LIGHT_CYAN = RGBColor(0xE0, 0xF7, 0xFA)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, corner_radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if corner_radius:
        shape.adjustments[0] = corner_radius
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False,
                 color=COLOR_DARK, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_paragraph(text_frame, text, font_size=16, bold=False, color=COLOR_DARK,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(6), space_after=Pt(2), bullet=False):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Microsoft YaHei"
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    if bullet:
        p.level = 0
    return p

def add_top_bar(slide):
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), COLOR_PRIMARY)

def add_page_number(slide, num, total):
    add_text_box(slide, Inches(12.2), Inches(7.1), Inches(1), Inches(0.3),
                 f"{num} / {total}", font_size=10, color=COLOR_GRAY, alignment=PP_ALIGN.RIGHT)

def add_section_header(slide, title, subtitle=""):
    add_top_bar(slide)
    shape = add_shape(slide, Inches(0.5), Inches(0.35), Inches(12.333), Inches(1.0),
                      COLOR_PRIMARY, corner_radius=0.02)
    add_text_box(slide, Inches(0.9), Inches(0.42), Inches(10), Inches(0.5),
                 title, font_size=32, bold=True, color=COLOR_WHITE)
    if subtitle:
        add_text_box(slide, Inches(0.9), Inches(0.88), Inches(10), Inches(0.35),
                     subtitle, font_size=15, color=RGBColor(0xBB, 0xDE, 0xFB))

def make_card(slide, left, top, width, height, bg_color, title, items, icon_text="",
              title_color=COLOR_DARK, item_color=COLOR_GRAY, title_size=18, item_size=14):
    card = add_shape(slide, left, top, width, height, bg_color, corner_radius=0.03)
    y_offset = top + Inches(0.15)
    if icon_text:
        add_text_box(slide, left + Inches(0.2), y_offset, Inches(0.4), Inches(0.4),
                     icon_text, font_size=22, bold=True, color=COLOR_PRIMARY)
        add_text_box(slide, left + Inches(0.65), y_offset, width - Inches(0.85), Inches(0.35),
                     title, font_size=title_size, bold=True, color=title_color)
    else:
        add_text_box(slide, left + Inches(0.2), y_offset, width - Inches(0.4), Inches(0.35),
                     title, font_size=title_size, bold=True, color=title_color)
    y_offset += Inches(0.45)
    txBox = slide.shapes.add_textbox(left + Inches(0.25), y_offset, width - Inches(0.5), height - Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(item_size)
        p.font.color.rgb = item_color
        p.font.name = "Microsoft YaHei"
        p.space_before = Pt(3)
        p.space_after = Pt(3)

TOTAL_SLIDES = 19

# ===================== Slide 1: 封面 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_WHITE)

add_shape(slide, Inches(0), Inches(0), Inches(5.5), Inches(7.5), COLOR_PRIMARY)
add_shape(slide, Inches(0), Inches(0), Inches(5.5), Inches(7.5), RGBColor(0x0D, 0x47, 0xA1))

for i in range(8):
    add_shape(slide, Inches(0.3 + i * 0.6), Inches(6.5), Inches(0.35), Inches(0.35),
              COLOR_ACCENT, corner_radius=0.5)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(4), Inches(1.2),
             "校园二手交易平台", font_size=40, bold=True, color=COLOR_WHITE)
add_text_box(slide, Inches(0.8), Inches(2.8), Inches(4), Inches(0.5),
             "需求分析报告", font_size=28, color=RGBColor(0xBB, 0xDE, 0xFB))

add_shape(slide, Inches(0.8), Inches(3.6), Inches(3), Inches(0.04), COLOR_ACCENT)

add_text_box(slide, Inches(0.8), Inches(4.2), Inches(4), Inches(0.8),
             "Campus Second-hand Trading Platform\nRequirements Analysis Report",
             font_size=14, color=RGBColor(0x90, 0xCA, 0xF9))

add_text_box(slide, Inches(0.8), Inches(5.5), Inches(4), Inches(1.0),
             "技术栈：Spring Boot + MyBatis-Plus + Vue 3 + Vite\n前端框架：Vue Router + Axios + Element Plus\n后端框架：JWT 认证 + Redis + WebSocket",
             font_size=13, color=RGBColor(0xBB, 0xDE, 0xFB))

add_text_box(slide, Inches(6.5), Inches(1.0), Inches(6), Inches(5.5), "", font_size=14)

make_card(slide, Inches(6.5), Inches(1.2), Inches(5.8), Inches(1.2), COLOR_LIGHT_BLUE,
          "项目定位", ["面向高校学生的校园二手物品交易平台，实现安全、便捷、可信赖的 C2C 交易体验"], icon_text="")

make_card(slide, Inches(6.5), Inches(2.7), Inches(5.8), Inches(1.2), COLOR_LIGHT_GREEN,
          "核心理念", ["信用体系 + 实名认证 + 担保交易 + 智能匹配，打造校园专属的可信交易环境"], icon_text="")

make_card(slide, Inches(6.5), Inches(4.2), Inches(5.8), Inches(1.2), COLOR_LIGHT_ORANGE,
          "覆盖范围", ["教材书籍、电子产品、生活用品、毕业甩卖等校园全品类二手交易场景"], icon_text="")

add_text_box(slide, Inches(7.5), Inches(6.2), Inches(4), Inches(0.5),
             "2026 年 4 月", font_size=14, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)

# ===================== Slide 2: 目录 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_BG)
add_top_bar(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
             "目  录", font_size=36, bold=True, color=COLOR_PRIMARY_DARK)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(5), Inches(0.4),
             "CONTENTS", font_size=16, color=COLOR_GRAY)
add_shape(slide, Inches(0.8), Inches(1.5), Inches(2), Inches(0.04), COLOR_PRIMARY)

toc_items = [
    ("01", "项目背景与目标", "分析校园二手交易痛点与平台建设目标"),
    ("02", "用户角色分析", "学生用户、管理员等多角色需求拆解"),
    ("03", "系统总体架构", "前后端分离架构与技术选型"),
    ("04", "核心功能模块概览", "平台六大核心功能模块总览"),
    ("05", "用户认证与信用体系", "实名认证、信用积分、等级机制"),
    ("06", "商品管理模块", "商品发布、审核、搜索、智能匹配"),
    ("07", "求购与匹配模块", "求购发布、供需智能匹配算法"),
    ("08", "订单与交易模块", "下单、担保交易、自提、确认收货"),
    ("09", "评价与纠纷模块", "互评机制、纠纷仲裁、恶意评价申诉"),
    ("10", "即时通讯与安全", "站内聊天、敏感词过滤、举报机制"),
    ("11", "校园特色功能", "毕业活动、教材匹配、校园公告"),
    ("12", "数据库表结构设计", "18张数据表的字段、约束与索引设计"),
    ("13", "数据库 ER 关系图", "核心实体关系与外键依赖可视化"),
    ("14", "非功能性需求", "性能、安全、可用性等质量需求"),
]

for i, (num, title, desc) in enumerate(toc_items):
    row = i // 3
    col = i % 3
    x = Inches(0.8 + col * 4.0)
    y = Inches(2.0 + row * 1.35)

    card_color = [COLOR_LIGHT_BLUE, COLOR_LIGHT_GREEN, COLOR_LIGHT_ORANGE,
                  COLOR_LIGHT_RED, COLOR_LIGHT_CYAN][i % 5]
    make_card(slide, x, y, Inches(3.7), Inches(1.15), card_color,
              f"{num}  {title}", [desc], title_size=15, item_size=12)

add_page_number(slide, 2, TOTAL_SLIDES)

# ===================== Slide 3: 项目背景与目标 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_BG)
add_section_header(slide, "01  项目背景与目标", "Background & Objectives")

make_card(slide, Inches(0.5), Inches(1.7), Inches(6), Inches(2.5), COLOR_WHITE,
          "校园二手交易现状痛点", [
              "  传统校园交易依赖 QQ 群、微信群、贴吧等非专业渠道，信息分散",
              "  缺乏信用保障机制，交易纠纷频发，买卖双方信任度低",
              "  没有统一的商品分类与搜索功能，供需匹配效率低下",
              "  线下交易缺乏安全保障，无担保机制和规范的自提流程",
              "  教材等高频交易品类缺乏标准化信息，重复发布低效",
          ], icon_text="", item_size=14)

make_card(slide, Inches(6.833), Inches(1.7), Inches(6), Inches(2.5), COLOR_WHITE,
          "平台建设目标", [
              "  构建一站式校园二手交易平台，覆盖发布-匹配-交易-售后全流程",
              "  建立基于实名认证的信用积分体系，提升交易可信度",
              "  实现供需智能匹配，提高商品与求购的对接效率",
              "  提供担保交易与规范自提流程，保障交易安全",
              "  集成即时通讯、评价纠纷、毕业活动等校园特色功能",
          ], icon_text="", item_size=14)

make_card(slide, Inches(0.5), Inches(4.5), Inches(3.8), Inches(2.6), COLOR_LIGHT_BLUE,
          "核心价值主张", [
              "安全可靠 — 实名认证 + 信用体系",
              "高效匹配 — 智能供需对接算法",
              "便捷交易 — 担保交易 + 校园自提点",
              "校园专属 — 教材/毕业/公告特色功能",
          ], icon_text="", title_size=16, item_size=13)

make_card(slide, Inches(4.633), Inches(4.5), Inches(3.8), Inches(2.6), COLOR_LIGHT_GREEN,
          "目标用户群体", [
              "在校大学生（主要用户群）",
              "毕业生（毕业季甩卖需求）",
              "新生（教材/生活用品购买需求）",
              "校园管理人员（平台监管）",
          ], icon_text="", title_size=16, item_size=13)

make_card(slide, Inches(8.766), Inches(4.5), Inches(3.8), Inches(2.6), COLOR_LIGHT_ORANGE,
          "预期成效", [
              "降低校园交易纠纷率 60%+",
              "提升供需匹配效率 3 倍+",
              "覆盖全校 80%+ 学生用户",
              "实现教材循环利用率提升",
          ], icon_text="", title_size=16, item_size=13)

add_page_number(slide, 3, TOTAL_SLIDES)

# ===================== Slide 4: 用户角色分析 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_BG)
add_section_header(slide, "02  用户角色分析", "User Role Analysis")

make_card(slide, Inches(0.5), Inches(1.7), Inches(3.9), Inches(5.4), COLOR_WHITE,
          "学生用户（买家/卖家）", [
              "【注册与认证】",
              "  注册账号、实名认证（学号/校园卡）",
              "  信用积分初始 100 分",
              "",
              "【商品交易】",
              "  发布/搜索/浏览二手商品",
              "  发布求购需求，等待智能匹配",
              "  下单购买、担保交易、线下自提",
              "",
              "【社交互动】",
              "  卖家买家即时通讯",
              "  交易完成互评（好评/中评/差评）",
              "  举报违规商品或用户",
              "",
              "【信用管理】",
              "  查看信用积分与等级",
              "  恶意评价申诉",
              "  封禁处罚申诉",
          ], icon_text="", title_size=16, item_size=12)

make_card(slide, Inches(4.7), Inches(1.7), Inches(3.9), Inches(5.4), COLOR_WHITE,
          "平台管理员", [
              "【内容审核】",
              "  审核商品发布（通过/驳回）",
              "  审核求购需求",
              "  审核实名认证申请",
              "",
              "【纠纷处理】",
              "  处理交易纠纷投诉",
              "  裁决处罚并扣除信用分",
              "  处理恶意评价申诉",
              "",
              "【违规管理】",
              "  处理用户举报",
              "  封禁/解封违规账号",
              "  记录违规处罚",
              "",
              "【平台运营】",
              "  发布校园公告通知",
              "  管理商品分类",
              "  管理自提点信息",
          ], icon_text="", title_size=16, item_size=12)

make_card(slide, Inches(8.9), Inches(1.7), Inches(3.9), Inches(5.4), COLOR_WHITE,
          "系统角色权限矩阵", [
              "功能              学生    管理员",
              "━━━━━━━━━━━━━━━━━━━━━━━━",
              "注册/登录          ✓        ✓",
              "实名认证           ✓        审核",
              "发布商品           ✓        审核",
              "发布求购           ✓        审核",
              "下单购买           ✓        —",
              "担保交易           ✓        管理",
              "即时通讯           ✓        —",
              "评价/申诉          ✓        处理",
              "举报               ✓        处理",
              "纠纷处理           提交     裁决",
              "商品分类管理       —        ✓",
              "自提点管理         —        ✓",
              "公告发布           —        ✓",
              "用户封禁           —        ✓",
              "━━━━━━━━━━━━━━━━━━━━━━━━",
              "信用积分           查询     管理",
          ], icon_text="", title_size=16, item_size=11)

add_page_number(slide, 4, TOTAL_SLIDES)

# ===================== Slide 5: 系统总体架构 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_BG)
add_section_header(slide, "03  系统总体架构", "System Architecture")

# 前端层
add_shape(slide, Inches(0.5), Inches(1.7), Inches(5.8), Inches(2.4), COLOR_LIGHT_BLUE, corner_radius=0.02)
add_text_box(slide, Inches(0.7), Inches(1.8), Inches(5.4), Inches(0.35),
             "前端展示层  Frontend", font_size=18, bold=True, color=COLOR_PRIMARY)
fe_items = [
    "Vue 3 + Vite 构建工具，组件化开发",
    "Vue Router 实现 SPA 路由管理",
    "Axios 封装 HTTP 请求，统一拦截器",
    "响应式 UI 设计，适配移动端浏览",
    "WebSocket 即时通讯客户端",
]
txBox = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(5.4), Inches(1.8))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(fe_items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"  {item}"
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_GRAY
    p.font.name = "Microsoft YaHei"
    p.space_before = Pt(3)

# 后端层
add_shape(slide, Inches(0.5), Inches(4.3), Inches(5.8), Inches(2.8), COLOR_LIGHT_GREEN, corner_radius=0.02)
add_text_box(slide, Inches(0.7), Inches(4.4), Inches(5.4), Inches(0.35),
             "后端服务层  Backend", font_size=18, bold=True, color=COLOR_GREEN)
be_items = [
    "Spring Boot 2.x 框架，RESTful API 设计",
    "MyBatis-Plus ORM，简化数据库操作",
    "JWT + Redis 实现无状态认证与 Token 管理",
    "WebSocket 实现服务端消息推送",
    "AOP 拦截器实现权限校验与日志记录",
    "敏感词过滤工具保护聊天安全",
]
txBox = slide.shapes.add_textbox(Inches(0.7), Inches(4.85), Inches(5.4), Inches(2.0))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(be_items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"  {item}"
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_GRAY
    p.font.name = "Microsoft YaHei"
    p.space_before = Pt(3)

# 数据层
add_shape(slide, Inches(6.6), Inches(1.7), Inches(6.2), Inches(5.4), COLOR_WHITE, corner_radius=0.02)
add_text_box(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(0.35),
             "数据与基础设施层  Data & Infrastructure", font_size=18, bold=True, color=COLOR_DARK)

data_sections = [
    ("MySQL 数据库", ["核心业务数据持久化", "17 张业务表，覆盖全场景", "MyBatis-Plus 自动填充/逻辑删除"]),
    ("Redis 缓存", ["Token 会话管理与黑名单", "热点数据缓存加速", "分布式锁支持"]),
    ("Elasticsearch", ["商品全文搜索引擎", "求购需求模糊匹配", "关键词高亮与排序"]),
    ("RabbitMQ 消息队列", ["异步消息通知", "订单状态变更事件", "匹配结果推送"]),
    ("文件存储", ["商品图片/证据图片上传", "校园卡照片存储", "10MB 单文件限制"]),
]

y = Inches(2.3)
for title, items in data_sections:
    add_text_box(slide, Inches(6.9), y, Inches(5.6), Inches(0.3),
                 title, font_size=14, bold=True, color=COLOR_PRIMARY)
    y += Inches(0.3)
    for item in items:
        add_text_box(slide, Inches(7.1), y, Inches(5.4), Inches(0.25),
                     f"  {item}", font_size=11, color=COLOR_GRAY)
        y += Inches(0.22)
    y += Inches(0.08)

# API 接口层标注
add_shape(slide, Inches(0.5), Inches(7.2), Inches(12.3), Inches(0.04), COLOR_PRIMARY)
add_text_box(slide, Inches(3.5), Inches(6.85), Inches(6), Inches(0.3),
             "RESTful API  |  JSON 数据交换  |  CORS 跨域支持  |  接口文档",
             font_size=12, color=COLOR_PRIMARY, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 5, TOTAL_SLIDES)

# ===================== Slide 6: 核心功能模块概览 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_BG)
add_section_header(slide, "04  核心功能模块概览", "Core Functional Modules Overview")

modules = [
    ("用户认证与信用", COLOR_LIGHT_BLUE, COLOR_PRIMARY, [
        "用户注册/登录（JWT 认证）",
        "实名认证（学号/校园卡）",
        "信用积分体系（100 分起）",
        "信用等级评定",
        "违规处罚与封禁机制",
    ]),
    ("商品管理", COLOR_LIGHT_GREEN, COLOR_GREEN, [
        "商品发布（图文/分类/定价）",
        "商品审核（管理员）",
        "商品搜索与浏览",
        "批量发布支持",
        "商品状态管理",
    ]),
    ("求购与智能匹配", COLOR_LIGHT_ORANGE, COLOR_ORANGE, [
        "求购需求发布",
        "关键词智能匹配",
        "匹配评分排序",
        "供需双向通知",
        "教材 ISBN 专项匹配",
    ]),
    ("订单与交易", COLOR_LIGHT_RED, COLOR_RED, [
        "下单购买流程",
        "担保交易机制",
        "线下自提点管理",
        "自提验证码确认",
        "自动确认收货（48h）",
    ]),
    ("评价与纠纷", COLOR_LIGHT_CYAN, COLOR_ACCENT, [
        "交易完成互评",
        "好评/中评/差评体系",
        "恶意评价申诉",
        "交易纠纷投诉",
        "管理员仲裁裁决",
    ]),
    ("即时通讯与安全", RGBColor(0xF3, 0xE5, 0xF5), RGBColor(0x9C, 0x27, 0xB0), [
        "买卖双方实时聊天",
        "WebSocket 消息推送",
        "敏感词过滤检测",
        "用户/商品举报",
        "违规记录追踪",
    ]),
]

for i, (title, bg, accent, items) in enumerate(modules):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.15)
    y = Inches(1.7 + row * 2.85)
    make_card(slide, x, y, Inches(3.85), Inches(2.6), bg, title, items, title_size=16, item_size=12)
    add_shape(slide, x + Inches(0.2), y + Inches(2.45), Inches(3.45), Inches(0.04), accent)

add_page_number(slide, 6, TOTAL_SLIDES)

# ===================== Slide 7: 用户认证与信用体系 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_BG)
add_section_header(slide, "05  用户认证与信用体系", "User Authentication & Credit System")

make_card(slide, Inches(0.5), Inches(1.7), Inches(6), Inches(2.6), COLOR_WHITE,
          "认证流程设计", [
              "1. 用户注册：用户名(3-20位) + 密码(6-20位) + 可选手机/邮箱/昵称",
              "2. 普通登录：用户名密码登录，JWT 签发 Token（有效期 24h）",
              "3. 实名认证：提交学号、姓名、院系、宿舍信息 + 校园卡照片",
              "4. 认证审核：管理员审核通过后 authStatus 变更为已认证",
              "5. CAS 对接：预留高校统一认证平台接口，可扩展单点登录",
          ], icon_text="", item_size=13)

make_card(slide, Inches(6.833), Inches(1.7), Inches(6), Inches(2.6), COLOR_WHITE,
          "JWT 认证机制", [
              "Token 签发：登录成功后生成 JWT，存储于 Redis（24h 过期）",
              "Token 校验：拦截器统一校验 Token 有效性与过期时间",
              "Token 黑名单：登出时加入黑名单，防止 Token 被复用",
              "接口白名单：注册/登录/商品列表等公开接口免认证",
              "Token 刷新：每次登录更新 Redis 中的 Token 记录",
          ], icon_text="", item_size=13)

make_card(slide, Inches(0.5), Inches(4.6), Inches(4), Inches(2.6), COLOR_LIGHT_BLUE,
          "信用积分规则", [
              "初始积分：100 分",
              "好评交易：+1 分",
              "中评交易：+0 分",
              "差评交易：-3 分",
              "违规举报属实：+1 分（月上限 3）",
              "捐赠物品：+2 分（月上限 6）",
              "30 天无违规：+2 分",
              "恶意评价：-3 分",
          ], icon_text="", item_size=12)

make_card(slide, Inches(4.8), Inches(4.6), Inches(3.9), Inches(2.6), COLOR_LIGHT_GREEN,
          "信用等级体系", [
              "EXCELLENT — 优秀（高信用）",
              "GOOD — 良好（默认）",
              "NORMAL — 一般",
              "POOR — 较差",
              "BLACKLIST — 黑名单",
              "",
              "等级影响：",
              "  交易权限 / 是否需担保",
          ], icon_text="", item_size=12)

make_card(slide, Inches(9.0), Inches(4.6), Inches(3.8), Inches(2.6), COLOR_LIGHT_RED,
          "违规处罚机制", [
              "封禁状态（临时/永久）",
              "封禁原因记录",
              "封禁时间管理",
              "信用分扣除",
              "违规记录存档",
              "",
              "申诉通道：",
              "  用户可提交封禁申诉",
          ], icon_text="", item_size=12)

add_page_number(slide, 7, TOTAL_SLIDES)

# ===================== Slide 8: 商品管理模块 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_BG)
add_section_header(slide, "06  商品管理模块", "Product Management Module")

make_card(slide, Inches(0.5), Inches(1.7), Inches(4), Inches(5.4), COLOR_WHITE,
          "商品发布与管理", [
              "【商品信息】",
              "  标题、描述、原价、售价",
              "  成色等级（新旧程度）",
              "  关键词标签（用于匹配）",
              "  多图上传（JSON 存储）",
              "  是否可议价标识",
              "",
              "【特色字段】",
              "  ISBN 编号（教材类）",
              "  自提位置与详情",
              "  批量发布支持（batchId）",
              "",
              "【状态管理】",
              "  在售 → 已预订 → 已售出",
              "  在售 → 已下架（卖家操作）",
              "  在售 → 审核驳回（管理员）",
          ], icon_text="", item_size=12)

make_card(slide, Inches(4.8), Inches(1.7), Inches(4), Inches(5.4), COLOR_WHITE,
          "商品审核与搜索", [
              "【审核流程】",
              "  发布后进入待审核状态",
              "  管理员审核通过/驳回",
              "  驳回时填写审核备注",
              "",
              "【搜索功能】",
              "  关键词模糊搜索",
              "  分类筛选（树形分类）",
              "  价格区间过滤",
              "  成色条件筛选",
              "  Elasticsearch 全文检索",
              "",
              "【浏览数据】",
              "  浏览次数统计",
              "  收藏/想要次数",
              "  匹配次数记录",
          ], icon_text="", item_size=12)

make_card(slide, Inches(9.1), Inches(1.7), Inches(3.7), Inches(5.4), COLOR_WHITE,
          "商品分类体系", [
              "树形分类结构",
              "parentId 支持父子关系",
              "排序序号控制显示",
              "分类图标配置",
              "分类状态启停",
              "",
              "示例分类：",
              "  教材书籍",
              "    ├── 计算机类",
              "    ├── 语言文学类",
              "    └── 理工类",
              "  电子产品",
              "    ├── 手机配件",
              "    ├── 电脑配件",
              "    └── 耳机音箱",
              "  生活用品",
              "  运动器材",
              "  服饰鞋帽",
          ], icon_text="", item_size=12)

add_page_number(slide, 8, TOTAL_SLIDES)

# ===================== Slide 9: 求购与匹配模块 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_BG)
add_section_header(slide, "07  求购与匹配模块", "Purchase Request & Smart Matching")

make_card(slide, Inches(0.5), Inches(1.7), Inches(6), Inches(2.5), COLOR_WHITE,
          "求购需求发布", [
              "买家可发布求购信息，包含：标题、描述、分类、预算区间（最低/最高价）",
              "成色要求、关键词标签、参考图片、自提位置偏好",
              "是否可议价、是否只接受高信用卖家、求购有效期设置",
              "发布后进入审核流程，审核通过后对卖家可见",
          ], icon_text="", item_size=13)

make_card(slide, Inches(6.833), Inches(1.7), Inches(6), Inches(2.5), COLOR_WHITE,
          "智能匹配算法", [
              "基于关键词的供需匹配：商品关键词 ↔ 求购关键词",
              "匹配评分机制：根据相关度计算 matchScore 并排序",
              "教材专项匹配：通过 ISBN 编号精确匹配教材供需",
              "匹配结果双向通知：通知卖家和买家，提升成交转化",
          ], icon_text="", item_size=13)

make_card(slide, Inches(0.5), Inches(4.5), Inches(6), Inches(2.7), COLOR_LIGHT_ORANGE,
          "匹配流程示意", [
              "  卖家发布商品 → 填写关键词/ISBN → 商品入库",
              "      ↓",
              "  买家发布求购 → 填写关键词/预算 → 求购入库",
              "      ↓",
              "  系统定时匹配 → 关键词比对 + ISBN 匹配 → 生成匹配记录",
              "      ↓",
              "  匹配结果推送 → 双方收到通知 → 买家查看商品 → 下单购买",
          ], icon_text="", item_size=13)

make_card(slide, Inches(6.833), Inches(4.5), Inches(6), Inches(2.7), COLOR_LIGHT_BLUE,
          "教材匹配特色功能", [
              "Textbook 数据库存储教材标准化信息：",
              "  ISBN、书名、作者、出版社、版次",
              "  对应课程名、适用专业、适用年级",
              "",
              "匹配逻辑：商品 ISBN = 求购 ISBN → 精确匹配",
              "辅助信息：课程名/专业/年级 → 提高匹配精准度",
          ], icon_text="", item_size=13)

add_page_number(slide, 9, TOTAL_SLIDES)

# ===================== Slide 10: 订单与交易模块 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_BG)
add_section_header(slide, "08  订单与交易模块", "Order & Transaction Module")

make_card(slide, Inches(0.5), Inches(1.7), Inches(4), Inches(5.4), COLOR_WHITE,
          "订单状态流转", [
              "【订单状态机】",
              "",
              "  待付款 (PENDING_PAYMENT)",
              "      ↓ 买家付款",
              "  待发货/待自提 (PAID)",
              "      ↓ 卖家确认发货",
              "  待收货 (SHIPPED)",
              "      ↓ 买家确认收货",
              "  已完成 (COMPLETED)",
              "",
              "  任意阶段 → 已取消 (CANCELLED)",
              "",
              "【自动确认】",
              "  超过 48 小时未确认 → 自动确认收货",
              "  可申请延长确认（+24h）",
          ], icon_text="", item_size=12)

make_card(slide, Inches(4.8), Inches(1.7), Inches(4), Inches(5.4), COLOR_WHITE,
          "担保交易机制", [
              "【触发条件】",
              "  金额超过阈值（50 元）自动启用",
              "  买家也可手动选择担保交易",
              "",
              "【担保流程】",
              "  1. 买家付款 → 资金由平台托管",
              "  2. 卖家发货 → 买家收到商品",
              "  3. 买家确认收货 → 资金释放给卖家",
              "",
              "【担保状态】",
              "  NONE → HELD → RELEASED",
              "  NONE → HELD → REFUNDED（纠纷退款）",
              "",
              "【安全保障】",
              "  确认截止时间管理",
              "  纠纷期间资金冻结",
          ], icon_text="", item_size=12)

make_card(slide, Inches(9.1), Inches(1.7), Inches(3.7), Inches(5.4), COLOR_WHITE,
          "线下自提流程", [
              "【自提点管理】",
              "  校园自提点信息",
              "  名称、地址、坐标",
              "  联系人、联系电话",
              "  营业时间",
              "",
              "【自提流程】",
              "  1. 买家下单时选择自提点",
              "  2. 预约自提时间",
              "  3. 系统生成自提验证码",
              "  4. 到达自提点出示验证码",
              "  5. 卖家验证后完成交接",
              "",
              "【安全机制】",
              "  验证码防伪",
              "  自提时间预约",
          ], icon_text="", item_size=12)

add_page_number(slide, 10, TOTAL_SLIDES)

# ===================== Slide 11: 评价与纠纷模块 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_BG)
add_section_header(slide, "09  评价与纠纷模块", "Review & Dispute Resolution")

make_card(slide, Inches(0.5), Inches(1.7), Inches(6), Inches(2.5), COLOR_WHITE,
          "评价体系设计", [
              "交易完成后双方互评：好评(+1信用)、中评(+0)、差评(-3信用)",
              "评价包含文字内容，支持回复功能",
              "评价时限：交易完成后 24 小时内必须评价",
              "恶意评价标记与申诉机制：被评方可提交申诉，管理员审核裁决",
          ], icon_text="", item_size=13)

make_card(slide, Inches(6.833), Inches(1.7), Inches(6), Inches(2.5), COLOR_WHITE,
          "纠纷处理流程", [
              "用户提交纠纷投诉 → 填写描述 + 上传证据图片",
              "管理员受理 → 审查证据 → 作出裁决",
              "裁决结果：退款/赔偿/驳回 + 信用分处罚",
              "纠纷处理时限：24 小时内处理，48 小时内可申诉",
          ], icon_text="", item_size=13)

make_card(slide, Inches(0.5), Inches(4.5), Inches(6), Inches(2.7), COLOR_LIGHT_RED,
          "纠纷状态流转", [
              "  待处理 (PENDING) → 用户提交投诉",
              "      ↓ 管理员受理",
              "  处理中 (PROCESSING) → 审查证据",
              "      ↓ 作出裁决",
              "  已解决 (RESOLVED) → 执行处罚/退款",
              "  或",
              "  已驳回 (REJECTED) → 投诉不成立",
          ], icon_text="", item_size=13)

make_card(slide, Inches(6.833), Inches(4.5), Inches(6), Inches(2.7), COLOR_LIGHT_ORANGE,
          "举报机制", [
              "举报目标类型：商品 / 用户 / 评价",
              "举报需填写原因并上传证据图片",
              "管理员审核举报 → 确认违规 → 执行处罚",
              "举报属实奖励：举报人 +1 信用分（月上限 3 次）",
              "处罚联动：违规记录 + 信用扣分 + 账号封禁",
          ], icon_text="", item_size=13)

add_page_number(slide, 11, TOTAL_SLIDES)

# ===================== Slide 12: 即时通讯与安全 + 校园特色 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_BG)
add_section_header(slide, "10  即时通讯与校园特色功能", "Chat, Security & Campus Features")

make_card(slide, Inches(0.5), Inches(1.7), Inches(4), Inches(2.5), COLOR_WHITE,
          "即时通讯功能", [
              "基于 WebSocket 的实时消息传输",
              "支持文本/图片等多种消息类型",
              "会话管理（conversationId）",
              "消息已读状态追踪",
              "关联订单上下文",
          ], icon_text="", item_size=13)

make_card(slide, Inches(4.8), Inches(1.7), Inches(4), Inches(2.5), COLOR_WHITE,
          "安全防护机制", [
              "敏感词过滤：聊天内容实时检测",
              "敏感消息标记（isSensitive）",
              "举报机制：商品/用户/评价举报",
              "违规记录追踪与处罚",
              "JWT Token 安全认证",
          ], icon_text="", item_size=13)

make_card(slide, Inches(9.1), Inches(1.7), Inches(3.7), Inches(2.5), COLOR_WHITE,
          "隐私保护设计", [
              "手机号脱敏显示",
              "真实姓名脱敏显示",
              "宿舍房间号脱敏显示",
              "密码字段永不返回",
              "个人信息按权限展示",
          ], icon_text="", item_size=13)

make_card(slide, Inches(0.5), Inches(4.5), Inches(4), Inches(2.7), COLOR_LIGHT_BLUE,
          "毕业季活动模块", [
              "发起毕业甩卖活动",
              "活动标题/描述/留言",
              "折扣类型配置",
              "  - 折扣比例",
              "  - 固定优惠金额",
              "活动时间区间设置",
              "活动状态管理",
          ], icon_text="", item_size=12)

make_card(slide, Inches(4.8), Inches(4.5), Inches(4), Inches(2.7), COLOR_LIGHT_GREEN,
          "校园公告系统", [
              "管理员发布公告通知",
              "公告标题/内容/类型",
              "置顶公告支持",
              "公告状态管理（发布/下架）",
              "关联发布管理员信息",
              "面向全校学生推送",
          ], icon_text="", item_size=12)

make_card(slide, Inches(9.1), Inches(4.5), Inches(3.7), Inches(2.7), COLOR_LIGHT_ORANGE,
          "自提点管理", [
              "校园自提点信息维护",
              "名称/地址/坐标",
              "校区归属标识",
              "联系人与电话",
              "营业时间配置",
              "自提点状态管理",
          ], icon_text="", item_size=12)

add_page_number(slide, 12, TOTAL_SLIDES)

# ===================== Slide 13: 数据库表结构总览 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_BG)
add_section_header(slide, "12  数据库表结构设计", "Database Table Structure Design")

db_tables = [
    ("user", "用户表", "id(主键), username(唯一), password, phone, email, avatar, nickname, student_no(唯一), real_name, department, dormitory_building, dormitory_room, campus_card_photo, auth_type, auth_status, credit_score, credit_level, status, ban_time, ban_reason, has_appealed_ban, last_login_time", "27字段"),
    ("product", "商品表", "id(主键), user_id→user, category_id→category, title, description, price, original_price, condition_level, keywords, isbn→textbook, images, is_negotiable, pickup_location, status, audit_status, view_count, want_count, match_count, is_batch, batch_id", "20字段"),
    ("orders", "订单表", "id(主键), order_no(唯一), product_id→product, buyer_id→user, seller_id→user, purchase_request_id→purchase_request, price, is_escrow, escrow_status, pickup_point_id→pickup_point, pickup_time, pickup_code, status, confirm_deadline, buyer_reviewed, seller_reviewed, review_deadline", "20字段"),
    ("purchase_request", "求购表", "id(主键), user_id→user, category_id→category, title, description, budget_min, budget_max, condition_requirement, keywords, reference_images, is_negotiable, only_high_credit, status, audit_status, match_count, expire_time", "20字段"),
    ("category", "商品分类表", "id(主键), name, parent_id(自引用), sort_order, icon, status", "6字段"),
    ("review", "评价表", "id(主键), order_id→orders, reviewer_id→user, reviewee_id→user, type, content, is_malicious, is_revoked, reply_content, appeal_status, appeal_reason, appeal_result", "16字段"),
    ("chat_message", "聊天消息表", "id(主键), conversation_id, sender_id→user, receiver_id→user, order_id→orders, content, msg_type, is_read, is_sensitive", "9字段"),
    ("dispute", "纠纷表", "id(主键), order_id→orders, complainant_id→user, respondent_id→user, type, description, evidence_images, status, result, penalty_score, penalty_user_id→user, handler_id→admin", "15字段"),
    ("report", "举报表", "id(主键), reporter_id→user, target_type, target_id, reason, evidence_images, status, handler_id→admin, handle_result, handle_time", "12字段"),
    ("credit_record", "信用记录表", "id(主键), user_id→user, order_id→orders, type, score_change, before_score, after_score, description, operator_id", "9字段"),
    ("match_record", "匹配记录表", "id(主键), product_id→product, purchase_request_id→purchase_request, match_score, seller_notified, buyer_notified, status", "7字段"),
    ("admin", "管理员表", "id(主键), username(唯一), password, real_name, role, campus, status", "7字段"),
    ("campus_notice", "校园通知表", "id(主键), admin_id→admin, title, content, type, is_top, status", "7字段"),
    ("textbook", "教材信息表", "id(主键), isbn(唯一), title, author, publisher, edition, course_name, major, grade", "9字段"),
    ("graduation_activity", "毕业活动表", "id(主键), user_id→user, title, description, message, discount_type, discount_value, start_time, end_time, status", "10字段"),
    ("pickup_point", "自提点表", "id(主键), name, address, longitude, latitude, campus, contact_person, contact_phone, open_time, status", "10字段"),
    ("violation_record", "违规记录表", "id(主键), user_id→user, type, target_type, target_id, description, penalty, score_deducted, handler_id→admin", "9字段"),
    ("operation_log", "操作日志表", "id(主键), operator_id, operator_type, module, action, target_type, target_id, detail, ip", "9字段"),
]

for i, (tname, comment, fields, count) in enumerate(db_tables):
    row = i // 3
    col = i % 3
    x = Inches(0.3 + col * 4.3)
    y = Inches(1.55 + row * 1.18)
    bg = [COLOR_LIGHT_BLUE, COLOR_LIGHT_GREEN, COLOR_LIGHT_ORANGE,
          COLOR_LIGHT_RED, COLOR_LIGHT_CYAN, RGBColor(0xF3, 0xE5, 0xF5)][i % 6]
    add_shape(slide, x, y, Inches(4.1), Inches(1.08), bg, corner_radius=0.02)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.02), Inches(3.2), Inches(0.22),
                 f"{tname}  ({comment})", font_size=9, bold=True, color=COLOR_PRIMARY)
    add_text_box(slide, x + Inches(3.3), y + Inches(0.02), Inches(0.7), Inches(0.22),
                 count, font_size=8, color=COLOR_GRAY, alignment=PP_ALIGN.RIGHT)
    txBox = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.25), Inches(3.9), Inches(0.78))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = fields
    p.font.size = Pt(7)
    p.font.color.rgb = COLOR_GRAY
    p.font.name = "Microsoft YaHei"
    p.space_before = Pt(0)
    p.space_after = Pt(0)

add_page_number(slide, 13, TOTAL_SLIDES)

# ===================== Slide 14: 数据库表详情 — 用户与商品核心表 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_BG)
add_section_header(slide, "12  数据库表结构详情 — 核心业务表", "Core Business Tables Detail")

def add_table_detail(slide, x, y, w, h, title, rows, bg_color=COLOR_WHITE, pk_icon="PK"):
    add_shape(slide, x, y, w, h, bg_color, corner_radius=0.02)
    add_shape(slide, x, y, w, Inches(0.3), COLOR_PRIMARY, corner_radius=0.02)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.02), w - Inches(0.2), Inches(0.25),
                 title, font_size=12, bold=True, color=COLOR_WHITE)
    txBox = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.35), w - Inches(0.2), h - Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, row_text in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = row_text
        p.font.size = Pt(8)
        p.font.color.rgb = COLOR_DARK
        p.font.name = "Consolas"
        p.space_before = Pt(1)
        p.space_after = Pt(1)

add_table_detail(slide, Inches(0.3), Inches(1.5), Inches(4.1), Inches(5.7),
    "user 用户表 (27字段)", [
        "PK  id              BIGINT       AUTO_INCREMENT",
        "UQ  username        VARCHAR(50)  NOT NULL",
        "    password        VARCHAR(255) BCrypt加密",
        "    phone           VARCHAR(20)  加密存储",
        "    email           VARCHAR(100)",
        "    avatar          VARCHAR(500)",
        "    nickname        VARCHAR(50)",
        "UQ  student_no      VARCHAR(50)  学号",
        "    real_name       VARCHAR(50)  加密存储",
        "    department      VARCHAR(100) 院系",
        "    dormitory_building VARCHAR(100)",
        "    dormitory_room  VARCHAR(50)  加密存储",
        "    campus_card_photo VARCHAR(500)",
        "    auth_type       TINYINT  0/1/2",
        "    auth_status     TINYINT  0/1/2/3",
        "    auth_time       DATETIME",
        "    credit_score    INT      默认100",
        "    credit_level    VARCHAR(20) 默认GOOD",
        "    status          TINYINT  0/1/2",
        "    ban_time        DATETIME",
        "    ban_reason      VARCHAR(500)",
        "    has_appealed_ban TINYINT  0/1",
        "    last_login_time DATETIME",
        "    created_time    DATETIME 自动填充",
        "    updated_time    DATETIME 自动更新",
        "    deleted         TINYINT  逻辑删除",
        "IDX: credit_score, auth_status",
    ])

add_table_detail(slide, Inches(4.6), Inches(1.5), Inches(4.1), Inches(5.7),
    "product 商品表 (20字段)", [
        "PK  id              BIGINT       AUTO_INCREMENT",
        "FK  user_id         BIGINT  → user.id",
        "FK  category_id     BIGINT  → category.id",
        "    title           VARCHAR(200) NOT NULL",
        "    description     TEXT",
        "    price           DECIMAL(10,2) NOT NULL",
        "    original_price  DECIMAL(10,2)",
        "    condition_level TINYINT  1-10",
        "    keywords        VARCHAR(500) 逗号分隔",
        "    isbn            VARCHAR(20) → textbook",
        "    images          TEXT     JSON数组",
        "    is_negotiable   TINYINT  0/1",
        "    pickup_location VARCHAR(200)",
        "    pickup_detail   VARCHAR(200) 加密",
        "    status          TINYINT  0/1/2/3/4",
        "    audit_status    TINYINT  0/1/2/3",
        "    audit_remark    VARCHAR(500)",
        "    audit_time      DATETIME",
        "    view_count      INT      默认0",
        "    want_count      INT      默认0",
        "    match_count     INT      默认0",
        "    is_batch        TINYINT  0/1",
        "    batch_id        VARCHAR(50)",
        "    created/updated/deleted 同上",
        "FT  ft_title_desc(title,description)",
        "IDX: user_id, category_id, status, price, isbn",
    ])

add_table_detail(slide, Inches(8.9), Inches(1.5), Inches(4.1), Inches(5.7),
    "orders 订单表 (20字段)", [
        "PK  id              BIGINT  AUTO_INCREMENT",
        "UQ  order_no        VARCHAR(50)  NOT NULL",
        "FK  product_id      BIGINT  → product.id",
        "FK  buyer_id        BIGINT  → user.id",
        "FK  seller_id       BIGINT  → user.id",
        "FK  purchase_request_id BIGINT→PR.id",
        "    price           DECIMAL(10,2) NOT NULL",
        "    is_escrow       TINYINT  0/1",
        "    escrow_status   TINYINT  0/1/2/3",
        "FK  pickup_point_id BIGINT  → PP.id",
        "    pickup_time     DATETIME",
        "    pickup_code     VARCHAR(20)",
        "    status          TINYINT  0-7",
        "    confirm_deadline DATETIME",
        "    buyer_reviewed  TINYINT  0/1",
        "    seller_reviewed TINYINT  0/1",
        "    review_deadline DATETIME",
        "    cancel_reason   VARCHAR(500)",
        "    remark          VARCHAR(500)",
        "    created/updated/deleted 同上",
        "",
        "IDX: buyer_id, seller_id, product_id, status",
    ])

add_page_number(slide, 14, TOTAL_SLIDES)

# ===================== Slide 15: 数据库表详情 — 求购/匹配/评价/聊天表 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_BG)
add_section_header(slide, "12  数据库表结构详情 — 交易与互动表", "Transaction & Interaction Tables")

add_table_detail(slide, Inches(0.3), Inches(1.5), Inches(3.15), Inches(5.7),
    "purchase_request 求购表 (20字段)", [
        "PK  id              BIGINT  AUTO_INCREMENT",
        "FK  user_id         BIGINT  → user.id",
        "FK  category_id     BIGINT  → category.id",
        "    title           VARCHAR(200) NOT NULL",
        "    description     TEXT",
        "    budget_min      DECIMAL(10,2)",
        "    budget_max      DECIMAL(10,2)",
        "    condition_requirement TINYINT",
        "    keywords        VARCHAR(500)",
        "    reference_images TEXT    JSON数组",
        "    is_negotiable   TINYINT  0/1",
        "    only_high_credit TINYINT 0/1",
        "    pickup_location VARCHAR(200)",
        "    status          TINYINT  0-5",
        "    audit_status    TINYINT  0-3",
        "    audit_remark    VARCHAR(500)",
        "    audit_time      DATETIME",
        "    match_count     INT      默认0",
        "    expire_time     DATETIME 30天过期",
        "    created/updated/deleted 同上",
        "IDX: user_id,category_id,status,expire_time",
    ])

add_table_detail(slide, Inches(3.6), Inches(1.5), Inches(3.15), Inches(5.7),
    "match_record 匹配表 (7字段)", [
        "PK  id              BIGINT  AUTO_INCREMENT",
        "FK  product_id      BIGINT  → product.id",
        "FK  purchase_request_id → PR.id",
        "    match_score     DECIMAL(5,2)",
        "    seller_notified TINYINT  0/1",
        "    buyer_notified  TINYINT  0/1",
        "    status          TINYINT  0/1",
        "    created_time    DATETIME",
        "",
        "UK(product_id,purchase_request_id)",
        "IDX: product_id, purchase_request_id",
        "",
        "",
        "",
        "",
        "",
        "",
        "",
        "",
        "",
        "",
    ])

add_table_detail(slide, Inches(6.9), Inches(1.5), Inches(3.05), Inches(5.7),
    "review 评价表 (16字段)", [
        "PK  id              BIGINT  AUTO_INCREMENT",
        "FK  order_id        BIGINT  → orders.id",
        "FK  reviewer_id     BIGINT  → user.id",
        "FK  reviewee_id     BIGINT  → user.id",
        "    type            TINYINT  1/2/3",
        "    content         VARCHAR(500)",
        "    is_malicious    TINYINT  0/1",
        "    is_revoked      TINYINT  0/1",
        "    reply_content   VARCHAR(500)",
        "    reply_time      DATETIME",
        "    appeal_status   TINYINT  0-3",
        "    appeal_reason   VARCHAR(500)",
        "    appeal_time     DATETIME",
        "    appeal_result   VARCHAR(500)",
        "    appeal_process_time DATETIME",
        "    is_visible      TINYINT  0/1",
        "    created/updated/deleted",
        "IDX: order_id,reviewer_id,reviewee_id",
    ])

add_table_detail(slide, Inches(10.1), Inches(1.5), Inches(3.0), Inches(5.7),
    "chat_message 聊天表 (9字段)", [
        "PK  id              BIGINT  AUTO_INCREMENT",
        "    conversation_id VARCHAR(100)",
        "FK  sender_id       BIGINT  → user.id",
        "FK  receiver_id     BIGINT  → user.id",
        "FK  order_id        BIGINT  → orders.id",
        "    content         TEXT     NOT NULL",
        "    msg_type        TINYINT  1/2/3",
        "    is_read         TINYINT  0/1",
        "    is_sensitive    TINYINT  0/1",
        "    created_time    DATETIME",
        "",
        "IDX: conversation_id, sender_id,",
        "     receiver_id, created_time",
        "",
        "",
        "",
        "",
        "",
        "",
        "",
    ])

add_page_number(slide, 15, TOTAL_SLIDES)

# ===================== Slide 16: 数据库表详情 — 管理/运营/辅助表 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_BG)
add_section_header(slide, "12  数据库表结构详情 — 管理与辅助表", "Admin & Auxiliary Tables")

add_table_detail(slide, Inches(0.3), Inches(1.5), Inches(2.55), Inches(2.7),
    "admin 管理员表 (7字段)", [
        "PK  id              BIGINT  AUTO_INCREMENT",
        "UQ  username        VARCHAR(50) NOT NULL",
        "    password        VARCHAR(255)",
        "    real_name       VARCHAR(50)",
        "    role            VARCHAR(30)",
        "    campus          VARCHAR(100)",
        "    status          TINYINT  0/1",
        "    created/updated/deleted",
    ])

add_table_detail(slide, Inches(3.0), Inches(1.5), Inches(2.55), Inches(2.7),
    "category 分类表 (6字段)", [
        "PK  id              BIGINT  AUTO_INCREMENT",
        "    name            VARCHAR(50) NOT NULL",
        "    parent_id       BIGINT  自引用",
        "    sort_order      INT      默认0",
        "    icon            VARCHAR(500)",
        "    status          TINYINT  0/1",
        "    created/updated/deleted",
        "树形结构: parent_id=0为顶级",
    ])

add_table_detail(slide, Inches(5.7), Inches(1.5), Inches(2.55), Inches(2.7),
    "textbook 教材表 (9字段)", [
        "PK  id              BIGINT  AUTO_INCREMENT",
        "UQ  isbn            VARCHAR(20) NOT NULL",
        "    title           VARCHAR(200) NOT NULL",
        "    author          VARCHAR(100)",
        "    publisher       VARCHAR(100)",
        "    edition         VARCHAR(50)",
        "    course_name     VARCHAR(100)",
        "    major           VARCHAR(100)",
        "    grade           VARCHAR(50)",
        "IDX: course_name, major",
    ])

add_table_detail(slide, Inches(8.4), Inches(1.5), Inches(2.3), Inches(2.7),
    "pickup_point 自提点表 (10字段)", [
        "PK  id              BIGINT  AUTO_INCREMENT",
        "    name            VARCHAR(100)",
        "    address         VARCHAR(300)",
        "    longitude       DECIMAL(10,6)",
        "    latitude        DECIMAL(10,6)",
        "    campus          VARCHAR(100)",
        "    contact_person  VARCHAR(50)",
        "    contact_phone   VARCHAR(20)",
        "    open_time       VARCHAR(100)",
        "    status          TINYINT  0/1",
    ])

add_table_detail(slide, Inches(10.85), Inches(1.5), Inches(2.3), Inches(2.7),
    "campus_notice 通知表 (7字段)", [
        "PK  id              BIGINT  AUTO_INCREMENT",
        "FK  admin_id        BIGINT  → admin.id",
        "    title           VARCHAR(200)",
        "    content         TEXT     NOT NULL",
        "    type            TINYINT  1-4",
        "    is_top          TINYINT  0/1",
        "    status          TINYINT  0/1",
        "    created/updated/deleted",
        "IDX: type, status",
    ])

add_table_detail(slide, Inches(0.3), Inches(4.5), Inches(2.55), Inches(2.7),
    "dispute 纠纷表 (15字段)", [
        "PK  id              BIGINT  AUTO_INCREMENT",
        "FK  order_id        BIGINT  → orders.id",
        "FK  complainant_id  BIGINT  → user.id",
        "FK  respondent_id   BIGINT  → user.id",
        "    type            TINYINT  1-4",
        "    description     TEXT     NOT NULL",
        "    evidence_images TEXT     JSON数组",
        "    status          TINYINT  0-5",
        "    result          VARCHAR(500)",
        "    penalty_score   INT",
        "FK  penalty_user_id BIGINT  → user.id",
        "FK  handler_id      BIGINT  → admin.id",
        "    handle_time     DATETIME",
        "    created/updated/deleted",
        "IDX: order_id,complainant_id,status",
    ])

add_table_detail(slide, Inches(3.0), Inches(4.5), Inches(2.55), Inches(2.7),
    "report 举报表 (12字段)", [
        "PK  id              BIGINT  AUTO_INCREMENT",
        "FK  reporter_id     BIGINT  → user.id",
        "    target_type     TINYINT  1-5",
        "    target_id       BIGINT",
        "    reason          VARCHAR(500)",
        "    evidence_images TEXT     JSON数组",
        "    status          TINYINT  0-3",
        "FK  handler_id      BIGINT  → admin.id",
        "    handle_result   VARCHAR(500)",
        "    handle_time     DATETIME",
        "    created/updated/deleted",
        "IDX: reporter_id, (type,target_id), status",
    ])

add_table_detail(slide, Inches(5.7), Inches(4.5), Inches(2.55), Inches(2.7),
    "credit_record 信用记录 (9字段)", [
        "PK  id              BIGINT  AUTO_INCREMENT",
        "FK  user_id         BIGINT  → user.id",
        "FK  order_id        BIGINT  → orders.id",
        "    type            VARCHAR(30)",
        "    score_change    INT  正/负",
        "    before_score    INT",
        "    after_score     INT",
        "    description     VARCHAR(500)",
        "    operator_id     BIGINT",
        "    created_time    DATETIME",
        "IDX: user_id, type, created_time",
    ])

add_table_detail(slide, Inches(8.4), Inches(4.5), Inches(2.3), Inches(2.7),
    "violation_record 违规 (9字段)", [
        "PK  id              BIGINT  AUTO_INCREMENT",
        "FK  user_id         BIGINT  → user.id",
        "    type            VARCHAR(50)",
        "    target_type     TINYINT  1-5",
        "    target_id       BIGINT",
        "    description     VARCHAR(500)",
        "    penalty         VARCHAR(200)",
        "    score_deducted  INT",
        "FK  handler_id      BIGINT  → admin.id",
        "    created_time    DATETIME",
        "IDX: user_id, type",
    ])

add_table_detail(slide, Inches(10.85), Inches(4.5), Inches(2.3), Inches(2.7),
    "graduation_activity 活动(10字段)", [
        "PK  id              BIGINT  AUTO_INCREMENT",
        "FK  user_id         BIGINT  → user.id",
        "    title           VARCHAR(200)",
        "    description     TEXT",
        "    message         TEXT  寄语",
        "    discount_type   TINYINT  1/2",
        "    discount_value  DECIMAL(5,2)",
        "    start_time      DATETIME",
        "    end_time        DATETIME",
        "    status          TINYINT  0/1",
        "    created/updated/deleted",
        "IDX: user_id, status",
    ])

add_page_number(slide, 16, TOTAL_SLIDES)

# ===================== Slide 17: 数据库操作日志表 + 表关系统计 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_BG)
add_section_header(slide, "12  数据库表结构详情 — 日志与关系统计", "Operation Log & Relationship Summary")

add_table_detail(slide, Inches(0.3), Inches(1.5), Inches(4.1), Inches(3.0),
    "operation_log 操作日志表 (9字段)", [
        "PK  id              BIGINT  AUTO_INCREMENT",
        "    operator_id     BIGINT",
        "    operator_type   TINYINT  1/2/3",
        "    module          VARCHAR(50) NOT NULL",
        "    action          VARCHAR(100) NOT NULL",
        "    target_type     VARCHAR(50)",
        "    target_id       BIGINT",
        "    detail          TEXT     JSON格式",
        "    ip              VARCHAR(50)",
        "    created_time    DATETIME",
        "",
        "IDX: (operator_id,operator_type),",
        "     module, created_time",
        "",
        "operator_type: 1-用户 2-管理员 3-系统",
    ])

make_card(slide, Inches(4.6), Inches(1.5), Inches(8.2), Inches(3.0), COLOR_WHITE,
          "表关系统计汇总", [
              "  共 18 张数据表，覆盖用户、商品、订单、评价、聊天、纠纷、管理等全部业务场景",
              "",
              "  外键关系总计 30+ 条：",
              "    user 表被引用 15+ 次（商品/订单/求购/评价/聊天/纠纷/举报/信用/违规/活动）",
              "    orders 表被引用 6 次（评价/纠纷/聊天/信用记录）",
              "    product 表被引用 2 次（订单/匹配记录）",
              "    purchase_request 表被引用 2 次（订单/匹配记录）",
              "    category 表被引用 2 次（商品/求购）",
              "    admin 表被引用 4 次（通知/纠纷/举报/违规记录处理人）",
              "",
              "  唯一约束 6 个：user.username, user.student_no, orders.order_no,",
              "    admin.username, textbook.isbn, match_record(product+PR)",
              "",
              "  全文索引 1 个：product 表 ft_title_desc(title, description)",
          ], icon_text="", item_size=11)

make_card(slide, Inches(0.3), Inches(4.8), Inches(6.3), Inches(2.4), COLOR_LIGHT_BLUE,
          "数据库设计特点", [
              "  InnoDB 引擎，支持事务与外键约束",
              "  utf8mb4 字符集，支持 emoji 等特殊字符",
              "  逻辑删除（deleted 字段），数据可恢复",
              "  自动时间戳（created_time / updated_time）",
              "  驼峰命名映射（MyBatis-Plus mapUnderscoreToCamelCase）",
              "  敏感数据加密存储（password/phone/real_name/dormitory_room）",
          ], icon_text="", item_size=11)

make_card(slide, Inches(6.8), Inches(4.8), Inches(6.0), Inches(2.4), COLOR_LIGHT_GREEN,
          "索引设计策略", [
              "  主键索引：所有表均使用 BIGINT AUTO_INCREMENT",
              "  唯一索引：username, student_no, order_no, isbn 等业务唯一字段",
              "  普通索引：user_id, category_id, status 等高频查询字段",
              "  组合索引：report(target_type, target_id), operation_log(operator)",
              "  全文索引：product(title, description) 支持中文全文检索",
          ], icon_text="", item_size=11)

add_page_number(slide, 17, TOTAL_SLIDES)

# ===================== Slide 18: ER 关系图 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_WHITE)
add_top_bar(slide)
add_text_box(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.45),
             "13  数据库 ER 关系图", font_size=28, bold=True, color=COLOR_PRIMARY_DARK)
add_text_box(slide, Inches(0.5), Inches(0.7), Inches(12), Inches(0.3),
             "Entity-Relationship Diagram — 核心实体关系可视化", font_size=14, color=COLOR_GRAY)

def draw_entity_box(slide, x, y, w, title, fields, color=COLOR_PRIMARY, bg=COLOR_WHITE):
    header_h = Inches(0.28)
    row_h = Inches(0.18)
    total_h = header_h + row_h * len(fields) + Inches(0.06)

    add_shape(slide, x, y, w, total_h, bg, corner_radius=0.01)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, header_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    add_text_box(slide, x + Inches(0.05), y + Inches(0.02), w - Inches(0.1), Inches(0.22),
                 title, font_size=8, bold=True, color=COLOR_WHITE)

    for i, field in enumerate(fields):
        fy = y + header_h + Inches(0.04) + row_h * i
        prefix = field[:2]
        fname = field[2:]
        prefix_color = COLOR_RED if prefix == "PK" else (COLOR_PRIMARY if prefix == "FK" else COLOR_GRAY)
        add_text_box(slide, x + Inches(0.04), fy, Inches(0.22), Inches(0.16),
                     prefix, font_size=6, bold=True, color=prefix_color)
        add_text_box(slide, x + Inches(0.26), fy, w - Inches(0.3), Inches(0.16),
                     fname, font_size=6.5, color=COLOR_DARK)
    return total_h

def draw_arrow(slide, x1, y1, x2, y2, color=COLOR_GRAY):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = Pt(1.2)

user_fields = ["PK id", "   username", "   password", "   phone/email", "   nickname", "   student_no", "   real_name", "   credit_score", "   credit_level", "   status"]
prod_fields = ["PK id", "FK user_id", "FK category_id", "   title", "   price", "   isbn", "   keywords", "   status", "   audit_status"]
order_fields = ["PK id", "   order_no", "FK product_id", "FK buyer_id", "FK seller_id", "   price", "   escrow_status", "   status"]
pr_fields = ["PK id", "FK user_id", "FK category_id", "   title", "   budget_min/max", "   keywords", "   status"]
review_fields = ["PK id", "FK order_id", "FK reviewer_id", "FK reviewee_id", "   type", "   appeal_status"]
chat_fields = ["PK id", "   conversation_id", "FK sender_id", "FK receiver_id", "   content", "   is_sensitive"]
dispute_fields = ["PK id", "FK order_id", "FK complainant_id", "FK respondent_id", "   status", "   penalty_score"]
report_fields = ["PK id", "FK reporter_id", "   target_type/id", "   reason", "   status"]
credit_fields = ["PK id", "FK user_id", "FK order_id", "   type", "   score_change"]
match_fields = ["PK id", "FK product_id", "FK purchase_request_id", "   match_score"]
cat_fields = ["PK id", "   name", "   parent_id", "   sort_order"]
admin_fields = ["PK id", "   username", "   role", "   campus", "   status"]
notice_fields = ["PK id", "FK admin_id", "   title", "   type", "   is_top"]
pp_fields = ["PK id", "   name", "   address", "   campus", "   contact_person"]
violation_fields = ["PK id", "FK user_id", "   type", "   penalty", "FK handler_id"]
grad_fields = ["PK id", "FK user_id", "   title", "   discount_type", "   status"]
textbook_fields = ["PK id", "   isbn", "   title", "   course_name", "   major"]
log_fields = ["PK id", "   operator_id", "   module", "   action", "   ip"]

erw = Inches(1.85)
erh_factor = 0.18

positions = {
    "user":     (Inches(4.5), Inches(1.1)),
    "product":  (Inches(7.0), Inches(1.1)),
    "category": (Inches(9.5), Inches(1.1)),
    "textbook": (Inches(11.2), Inches(1.1)),
    "orders":   (Inches(7.0), Inches(3.5)),
    "pr":       (Inches(4.5), Inches(3.5)),
    "match":    (Inches(3.0), Inches(5.2)),
    "review":   (Inches(9.5), Inches(3.5)),
    "chat":     (Inches(11.2), Inches(3.5)),
    "credit":   (Inches(1.0), Inches(3.5)),
    "dispute":  (Inches(9.5), Inches(5.5)),
    "report":   (Inches(7.0), Inches(5.5)),
    "admin":    (Inches(1.0), Inches(1.1)),
    "notice":   (Inches(1.0), Inches(5.5)),
    "pp":       (Inches(4.5), Inches(5.5)),
    "violation":(Inches(1.0), Inches(6.0)),
    "grad":     (Inches(3.0), Inches(6.5)),
    "log":      (Inches(11.2), Inches(5.5)),
}

entity_defs = [
    ("user", "user 用户", user_fields, COLOR_PRIMARY),
    ("product", "product 商品", prod_fields, COLOR_GREEN),
    ("category", "category 分类", cat_fields, COLOR_ACCENT),
    ("orders", "orders 订单", order_fields, COLOR_RED),
    ("pr", "purchase_request 求购", pr_fields, COLOR_ORANGE),
    ("review", "review 评价", review_fields, RGBColor(0x9C, 0x27, 0xB0)),
    ("chat", "chat_message 聊天", chat_fields, RGBColor(0x00, 0x96, 0x88)),
    ("dispute", "dispute 纠纷", dispute_fields, COLOR_RED),
    ("report", "report 举报", report_fields, RGBColor(0x79, 0x55, 0x48)),
    ("credit", "credit_record 信用", credit_fields, COLOR_PRIMARY_DARK),
    ("match", "match_record 匹配", match_fields, COLOR_ACCENT),
    ("admin", "admin 管理员", admin_fields, RGBColor(0x45, 0x5A, 0x64)),
    ("notice", "campus_notice 通知", notice_fields, RGBColor(0x00, 0x96, 0x88)),
    ("pp", "pickup_point 自提点", pp_fields, COLOR_GREEN),
    ("textbook", "textbook 教材", textbook_fields, COLOR_ORANGE),
    ("violation", "violation_record 违规", violation_fields, COLOR_RED),
    ("grad", "graduation_activity 活动", grad_fields, RGBColor(0x9C, 0x27, 0xB0)),
    ("log", "operation_log 日志", log_fields, RGBColor(0x45, 0x5A, 0x64)),
]

for key, title, fields, color in entity_defs:
    px, py = positions[key]
    draw_entity_box(slide, px, py, erw, title, fields, color=color)

# Draw relationship lines
rel_lines = [
    ("user", "product", 0, 1), ("category", "product", 0, 2),
    ("user", "orders", 0, 3), ("product", "orders", 0, 2),
    ("user", "pr", 0, 1), ("category", "pr", 0, 2),
    ("product", "match", 0, 1), ("pr", "match", 0, 2),
    ("orders", "review", 0, 1), ("user", "review", 0, 2),
    ("user", "chat", 0, 2), ("user", "chat", 0, 3),
    ("orders", "dispute", 0, 1), ("user", "dispute", 0, 2),
    ("user", "report", 0, 1), ("user", "credit", 0, 1),
    ("orders", "credit", 0, 2), ("admin", "notice", 0, 1),
    ("user", "grad", 0, 1), ("user", "violation", 0, 1),
    ("admin", "dispute", 0, 4), ("admin", "report", 0, 3),
    ("admin", "violation", 0, 4), ("user", "orders", 0, 4),
    ("pp", "orders", 0, 4), ("orders", "chat", 0, 4),
]

drawn = set()
for src, dst, _, _ in rel_lines:
    key = (src, dst)
    if key in drawn:
        continue
    drawn.add(key)
    sx, sy = positions[src]
    dx, dy = positions[dst]
    src_h = Inches(0.28) + Inches(0.18) * len([f for f in [user_fields, prod_fields, cat_fields, order_fields, pr_fields, review_fields, chat_fields, dispute_fields, report_fields, credit_fields, match_fields, admin_fields, notice_fields, pp_fields, textbook_fields, violation_fields, grad_fields, log_fields][[e[0] for e in entity_defs].index(src)] if f.strip()])
    dst_h = Inches(0.28) + Inches(0.18) * len([f for f in [user_fields, prod_fields, cat_fields, order_fields, pr_fields, review_fields, chat_fields, dispute_fields, report_fields, credit_fields, match_fields, admin_fields, notice_fields, pp_fields, textbook_fields, violation_fields, grad_fields, log_fields][[e[0] for e in entity_defs].index(dst)] if f.strip()])
    draw_arrow(slide, sx + erw // 2, sy + src_h, dx + erw // 2, dy, COLOR_GRAY)

add_page_number(slide, 18, TOTAL_SLIDES)

# ===================== Slide 14: 非功能性需求 =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, COLOR_BG)
add_section_header(slide, "12  非功能性需求", "Non-functional Requirements")

make_card(slide, Inches(0.5), Inches(1.7), Inches(3.9), Inches(2.5), COLOR_LIGHT_BLUE,
          "性能需求", [
              "页面加载时间 < 3 秒",
              "API 响应时间 < 500ms",
              "支持 200+ 并发用户",
              "数据库连接池：最大 20 连接",
              "Redis 缓存热点数据",
              "ES 全文检索加速搜索",
          ], icon_text="", item_size=12)

make_card(slide, Inches(4.7), Inches(1.7), Inches(3.9), Inches(2.5), COLOR_LIGHT_GREEN,
          "安全需求", [
              "JWT 无状态认证 + Token 黑名单",
              "密码 BCrypt 加密存储",
              "CORS 跨域安全配置",
              "敏感词实时过滤",
              "个人信息脱敏展示",
              "文件上传大小限制（10MB）",
          ], icon_text="", item_size=12)

make_card(slide, Inches(8.9), Inches(1.7), Inches(3.9), Inches(2.5), COLOR_LIGHT_ORANGE,
          "可用性需求", [
              "前后端分离，独立部署",
              "RESTful API 标准化设计",
              "响应式 UI，适配多终端",
              "友好的错误提示信息",
              "操作日志可追溯",
              "数据逻辑删除，可恢复",
          ], icon_text="", item_size=12)

make_card(slide, Inches(0.5), Inches(4.5), Inches(3.9), Inches(2.7), COLOR_WHITE,
          "可维护性需求", [
              "分层架构：Controller → Service → Mapper",
              "MyBatis-Plus 简化 CRUD 操作",
              "统一异常处理（GlobalExceptionHandler）",
              "统一响应封装（Result 工具类）",
              "代码规范：驼峰映射、自动填充",
          ], icon_text="", item_size=12)

make_card(slide, Inches(4.7), Inches(4.5), Inches(3.9), Inches(2.7), COLOR_WHITE,
          "可扩展性需求", [
              "模块化设计，功能解耦",
              "预留 CAS 统一认证接口",
              "可扩展支付渠道对接",
              "可扩展推送通知渠道",
              "数据表设计支持功能迭代",
          ], icon_text="", item_size=12)

make_card(slide, Inches(8.9), Inches(4.5), Inches(3.9), Inches(2.7), COLOR_WHITE,
          "技术栈总结", [
              "后端：Spring Boot + MyBatis-Plus",
              "前端：Vue 3 + Vite + Vue Router",
              "数据库：MySQL + Redis + ES",
              "消息队列：RabbitMQ",
              "认证：JWT + Redis",
              "通讯：WebSocket + REST API",
          ], icon_text="", item_size=12)

add_page_number(slide, 19, TOTAL_SLIDES)

output_path = r"c:\Users\周云富\Desktop\校园二手交易平台\校园二手交易平台_需求分析报告.pptx"
prs.save(output_path)
print(f"PPT 已生成：{output_path}")
